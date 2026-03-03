"""
brd_parser.py
Extract plain text from uploaded BRD files (PDF / DOCX / PPTX).
"""
from __future__ import annotations

import io
from typing import Union


def parse_pdf(file_bytes: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    parts: list[str] = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            parts.append(text.strip())
    return "\n\n".join(parts)


def parse_docx(file_bytes: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    parts: list[str] = []

    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)

    return "\n".join(parts)


def parse_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)
        if slide_parts:
            parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_parts))

    return "\n\n".join(parts)


def parse_brd(filename: str, file_bytes: bytes) -> str:
    """Route to correct parser based on file extension."""
    ext = filename.rsplit(".", 1)[-1].lower()
    if ext == "pdf":
        return parse_pdf(file_bytes)
    elif ext in ("docx", "doc"):
        return parse_docx(file_bytes)
    elif ext in ("pptx", "ppt"):
        return parse_pptx(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: .{ext}. Please upload PDF, DOCX, or PPTX.")
